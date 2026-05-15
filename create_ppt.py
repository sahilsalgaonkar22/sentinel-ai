"""
Sentinel AI — Professional 7-Slide Presentation Generator
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color Palette ───
BG_DARK    = RGBColor(0x08, 0x0C, 0x1A)
BG_CARD    = RGBColor(0x0F, 0x16, 0x29)
PURPLE     = RGBColor(0xBD, 0x9D, 0xFF)
CYAN       = RGBColor(0x36, 0xBC, 0xFD)
GREEN      = RGBColor(0x2E, 0xD5, 0x73)
ORANGE     = RGBColor(0xFF, 0x6B, 0x35)
RED        = RGBColor(0xEF, 0x44, 0x44)
WHITE      = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY  = RGBColor(0x64, 0x74, 0x8B)
YELLOW     = RGBColor(0xEA, 0xB3, 0x08)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, bullet_color=PURPLE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE
    shape.line.fill.background()
    shape.shadow.inherit = False

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)

# Accent line at top
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), PURPLE)

# Title
add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
         "SENTINEL AI", font_size=54, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
         "Sovereign Observer System", font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Tagline
add_text(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.8),
         "AI-Powered Enterprise Security Intelligence Platform\nfor Automated Vulnerability Detection, Risk Analysis & Compliance",
         font_size=18, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom info
add_accent_line(slide, Inches(5.5), Inches(5.5), Inches(2.3))
add_text(slide, Inches(1.5), Inches(6), Inches(10), Inches(0.5),
         "National Forensic Sciences University  •  M.Tech Cyber Security  •  2026",
         font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(6.4), Inches(10), Inches(0.5),
         "Presented by: Sahil",
         font_size=16, color=LIGHT_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement & Objectives
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), PURPLE)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
         "Problem Statement & Objectives", font_size=30, color=PURPLE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

# Problem box
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), BG_CARD, RGBColor(0x1E, 0x3A, 0x5F))
add_text(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5),
         "⚠  The Problem", font_size=20, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(3.5), [
    "▸  Organizations face increasing cyber threats but lack\n   centralized, AI-driven security platforms",
    "▸  Manual vulnerability assessment is slow, error-prone,\n   and fails to scale across infrastructure",
    "▸  No unified system to scan, detect, analyze, and\n   report in a single pipeline",
    "▸  Compliance reporting (OWASP, PCI-DSS, ISO 27001)\n   is done manually — costly and incomplete",
], font_size=14, color=LIGHT_GRAY)

# Objectives box
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5), BG_CARD, RGBColor(0x1E, 0x3A, 0x5F))
add_text(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5),
         "🎯  Objectives", font_size=20, color=GREEN, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(2.3), Inches(5.2), Inches(3.5), [
    "▸  Build an AI-powered security platform that automates\n   end-to-end vulnerability management",
    "▸  Integrate 8+ real security scanners (Nmap, Nikto,\n   Semgrep, Trivy, Nuclei, Bandit, etc.)",
    "▸  Implement real-time threat intelligence with EPSS\n   and CISA KEV enrichment",
    "▸  Auto-generate compliance reports mapped to OWASP\n   Top 10, PCI-DSS, ISO 27001, NIST CSF",
    "▸  Deliver a production-grade, containerized platform\n   with enterprise-class architecture",
], font_size=14, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — System Architecture
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), PURPLE)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
         "System Architecture", font_size=30, color=PURPLE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

# Architecture layers - horizontal flow
layers = [
    ("Frontend\n(React + Vite)", "React 18, Framer Motion\nPremium Dark UI\nReal-time WebSocket", CYAN),
    ("API Gateway\n(FastAPI)", "RESTful + WebSocket\nJWT Authentication\nRate Limiting + CORS", PURPLE),
    ("Event Bus\n(Apache Kafka)", "Async Scan Dispatch\nDead Letter Queue\nResult Processing", ORANGE),
    ("Worker Fleet\n(Distributed)", "8+ Security Scanners\nParallel Execution\nThreat Enrichment", GREEN),
    ("Data Layer\n(Multi-Store)", "PostgreSQL + Redis\nElasticsearch + MinIO\nGrafana + Prometheus", YELLOW),
]

for i, (title, desc, color) in enumerate(layers):
    x = Inches(0.6 + i * 2.5)
    y = Inches(1.6)
    card = add_shape(slide, x, y, Inches(2.2), Inches(2.8), BG_CARD, color)
    add_text(slide, x + Inches(0.15), y + Inches(0.2), Inches(1.9), Inches(0.7),
             title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.15), y + Inches(1.0), Inches(1.9), Inches(1.5),
             desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < len(layers) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.2), y + Inches(1.2), Inches(0.3), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_GRAY
        arrow.line.fill.background()

# Workflow footer
add_shape(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2), BG_CARD, RGBColor(0x1E, 0x3A, 0x5F))
add_text(slide, Inches(0.9), Inches(4.95), Inches(5), Inches(0.4),
         "🔄  End-to-End Scan Workflow", font_size=16, color=PURPLE, bold=True)
add_text(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.2),
         "User submits target  →  Input detector identifies target type  →  Orchestrator selects tools  →  "
         "Kafka dispatches to worker fleet  →  Workers execute scans in parallel  →  Results enriched with EPSS/KEV  →  "
         "Findings stored in PostgreSQL  →  AI engine computes risk score  →  PDF report auto-generated with compliance mapping  →  "
         "Dashboard updated in real-time via WebSocket",
         font_size=12, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Technology Stack & Security Tools
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), PURPLE)

add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
         "Technology Stack & Security Tools", font_size=30, color=PURPLE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

# Tech stack table
categories = [
    ("Frontend", "React 18, Vite, Framer Motion, Recharts, Lucide Icons", CYAN),
    ("Backend", "Python, FastAPI, SQLAlchemy, Uvicorn, Async/Await", PURPLE),
    ("Messaging", "Apache Kafka with DLQ, Redis for caching & rate limiting", ORANGE),
    ("Database", "PostgreSQL (primary), Elasticsearch (search), MinIO (objects)", GREEN),
    ("Monitoring", "Prometheus metrics, Grafana dashboards, structured JSON logging", YELLOW),
    ("Auth", "JWT tokens, bcrypt hashing, org-scoped multi-tenancy, RBAC", RED),
    ("Containerization", "Docker Compose orchestration, 10 microservices, health checks", CYAN),
]

for i, (cat, desc, color) in enumerate(categories):
    y = Inches(1.5 + i * 0.7)
    add_shape(slide, Inches(0.8), y, Inches(2.2), Inches(0.55), BG_CARD, color)
    add_text(slide, Inches(0.95), y + Inches(0.08), Inches(2), Inches(0.4),
             cat, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.2), y + Inches(0.08), Inches(5), Inches(0.4),
             desc, font_size=13, color=LIGHT_GRAY)

# Security tools
add_shape(slide, Inches(8.8), Inches(1.5), Inches(3.8), Inches(5.2), BG_CARD, RGBColor(0x1E, 0x3A, 0x5F))
add_text(slide, Inches(9.0), Inches(1.65), Inches(3.5), Inches(0.4),
         "🛡  Integrated Security Scanners", font_size=14, color=PURPLE, bold=True)

tools = [
    ("Nmap", "Network discovery & port scanning"),
    ("Nikto", "Web server vulnerability scanning"),
    ("Semgrep", "Static code analysis (SAST)"),
    ("Bandit", "Python security linter"),
    ("Trivy", "Container & dependency scanning"),
    ("Nuclei", "Template-based vuln detection"),
    ("Subfinder", "Subdomain enumeration"),
    ("httpx", "HTTP probing & tech detection"),
    ("Gitleaks", "Secret & credential detection"),
]

for i, (tool, desc) in enumerate(tools):
    y = Inches(2.2 + i * 0.48)
    add_text(slide, Inches(9.1), y, Inches(1.2), Inches(0.35),
             f"▸ {tool}", font_size=12, color=GREEN, bold=True)
    add_text(slide, Inches(10.2), y, Inches(2.3), Inches(0.35),
             desc, font_size=10, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Key Features
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), PURPLE)

add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
         "Key Features", font_size=30, color=PURPLE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

features = [
    ("📊", "Real-Time Dashboard", "Command Center with live risk index,\nthreat stream, and system health monitoring", CYAN),
    ("🔍", "Multi-Tool Scanning", "9 integrated scanners with auto-detection\nof target type and parallel execution", GREEN),
    ("🧠", "AI Risk Engine", "CVSS + EPSS + CISA KEV enrichment for\nintelligent risk scoring and prioritization", PURPLE),
    ("📄", "PDF Report Generation", "One-click executive reports with findings,\ncharts, severity breakdown, compliance data", ORANGE),
    ("🛡", "Compliance Mapping", "Auto-map findings to OWASP Top 10, PCI-DSS\nv4.0, ISO 27001:2022, NIST CSF 2.0", RED),
    ("📈", "Scan Drift Detection", "Compare scans over time — track new vulns,\nresolved issues, and security posture change", YELLOW),
    ("⏱", "Scheduled Scans", "Cron-based recurring scans with automated\nalert triggers on critical findings", CYAN),
    ("🔔", "Slack / Email Alerts", "Real-time notifications via SMTP and Slack\nwebhooks with Redis-backed rate limiting", GREEN),
]

for i, (icon, title, desc, color) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.6 + col * 3.1)
    y = Inches(1.5 + row * 2.7)
    card = add_shape(slide, x, y, Inches(2.9), Inches(2.3), BG_CARD, color)
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.5),
             f"{icon}  {title}", font_size=16, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.5), Inches(1.2),
             desc, font_size=12, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Results & Compliance
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), PURPLE)

add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
         "Results & Compliance Validation", font_size=30, color=PURPLE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

# Test results box
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), BG_CARD, RGBColor(0x1E, 0x3A, 0x5F))
add_text(slide, Inches(1.1), Inches(1.65), Inches(5), Inches(0.4),
         "✅  System Audit Results", font_size=18, color=GREEN, bold=True)

results = [
    "▸  43 / 43 API endpoint tests — PASSED",
    "▸  10 / 10 Docker containers — HEALTHY",
    "▸  6 / 6 enterprise features — VERIFIED",
    "▸  Authentication & multi-tenancy — WORKING",
    "▸  Real scan execution on live targets — CONFIRMED",
    "▸  PDF report generation pipeline — OPERATIONAL",
    "▸  Kafka event processing — 0 message loss",
    "▸  WebSocket real-time updates — ACTIVE",
    "▸  Zero mock/simulated data — ALL REAL",
]
add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(4), results, font_size=14, color=LIGHT_GRAY)

# Compliance box
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2), BG_CARD, RGBColor(0x1E, 0x3A, 0x5F))
add_text(slide, Inches(7.1), Inches(1.65), Inches(5), Inches(0.4),
         "🛡  Compliance Framework Mapping", font_size=18, color=PURPLE, bold=True)

# Framework cards
frameworks = [
    ("OWASP Top 10 (2021)", "8 / 10 controls passed", "2 violations detected", ORANGE),
    ("PCI-DSS v4.0", "10 / 11 controls passed", "1 violation detected", RED),
    ("ISO 27001:2022", "11 / 11 controls passed", "Fully Compliant ✓", GREEN),
    ("NIST CSF 2.0", "7 / 8 controls passed", "1 gap identified", CYAN),
]

for i, (name, score, status, color) in enumerate(frameworks):
    y = Inches(2.3 + i * 1.1)
    add_shape(slide, Inches(7.1), y, Inches(5.1), Inches(0.9), RGBColor(0x0A, 0x0F, 0x1E), color)
    add_text(slide, Inches(7.3), y + Inches(0.08), Inches(2.5), Inches(0.35),
             name, font_size=13, color=color, bold=True)
    add_text(slide, Inches(9.8), y + Inches(0.08), Inches(1.5), Inches(0.35),
             score, font_size=12, color=WHITE, bold=True)
    add_text(slide, Inches(7.3), y + Inches(0.45), Inches(4.5), Inches(0.35),
             status, font_size=11, color=LIGHT_GRAY)

# Overall score
add_text(slide, Inches(7.1), Inches(6.0), Inches(5.1), Inches(0.5),
         "Overall Compliance Score:  90%  —  PARTIAL COMPLIANCE", font_size=16, color=YELLOW, bold=True,
         alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Conclusion & Future Scope
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), PURPLE)

add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
         "Conclusion & Future Scope", font_size=30, color=PURPLE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.8))

# Conclusion
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5), BG_CARD, RGBColor(0x1E, 0x3A, 0x5F))
add_text(slide, Inches(1.1), Inches(1.65), Inches(5), Inches(0.4),
         "📌  Conclusion", font_size=20, color=GREEN, bold=True)
conclusions = [
    "▸  Successfully built a production-grade, AI-powered\n   security intelligence platform",
    "▸  Demonstrated real vulnerability detection on live\n   targets with 9 integrated scanning tools",
    "▸  Achieved 90% automated compliance mapping across\n   4 major security frameworks",
    "▸  Delivered enterprise features: multi-tenancy, PDF\n   reports, scheduled scans, alert notifications",
    "▸  Platform is containerized with Docker Compose for\n   seamless deployment and scalability",
]
add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(3.5), conclusions, font_size=14, color=LIGHT_GRAY)

# Future scope
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.5), BG_CARD, RGBColor(0x1E, 0x3A, 0x5F))
add_text(slide, Inches(7.1), Inches(1.65), Inches(5), Inches(0.4),
         "🚀  Future Scope", font_size=20, color=CYAN, bold=True)
future = [
    "▸  LLM-powered remediation suggestions using\n   GPT/Gemini for finding-specific fix guidance",
    "▸  Kubernetes-native deployment with auto-scaling\n   worker pods based on scan queue depth",
    "▸  SIEM integration (Splunk, QRadar) for enterprise\n   SOC workflow automation",
    "▸  Mobile companion app with real-time push\n   notifications for critical findings",
    "▸  Threat intelligence feed integration (MITRE ATT&CK,\n   AlienVault OTX, VirusTotal)",
]
add_bullet_list(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(3.5), future, font_size=14, color=LIGHT_GRAY)

# Thank you footer
add_shape(slide, Inches(3), Inches(6.3), Inches(7.3), Inches(0.8), BG_CARD, PURPLE)
add_text(slide, Inches(3.2), Inches(6.4), Inches(6.9), Inches(0.6),
         "Thank You  —  Questions?", font_size=24, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = r"C:\Users\SAHIL\Downloads\ransom2\Sentinel_AI_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
